"""Universal ingestion: any source → semantic docs (+ procedure candidates).

Format handlers degrade gracefully: if an optional dependency is missing the
handler reports what to install instead of crashing the pipeline.
"""

from __future__ import annotations

import uuid
from collections.abc import Callable
from pathlib import Path

from atlas.ingest.procedures import extract_procedure
from atlas.memory.semantic import Doc, SemanticStore

_Handler = Callable[[Path], str]
_HANDLERS: dict[str, _Handler] = {}


def _register(*exts: str) -> Callable[[_Handler], _Handler]:
    def deco(fn: _Handler) -> _Handler:
        for ext in exts:
            _HANDLERS[ext] = fn
        return fn
    return deco


@_register(".pdf")
def _pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "[install pypdf to ingest PDFs: pip install pypdf]"
    return "\n".join(page.extract_text() or "" for page in PdfReader(str(path)).pages)


@_register(".pptx")
def _pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "[install python-pptx to ingest PPTX]"
    texts = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


@_register(".docx")
def _docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        return "[install python-docx to ingest DOCX]"
    return "\n".join(p.text for p in Document(str(path)).paragraphs)


@_register(".txt", ".md", ".srt", ".vtt")
def _text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def fetch_video_transcript(url: str, workdir: str) -> Path | None:
    """Download subtitles for a video URL via yt-dlp (auto-subs fallback)."""
    import subprocess

    out = Path(workdir)
    out.mkdir(parents=True, exist_ok=True)
    cmd = [
        "yt-dlp", "--skip-download", "--write-subs", "--write-auto-subs",
        "--sub-langs", "en.*", "--convert-subs", "srt",
        "-o", str(out / "%(id)s"), url,
    ]
    try:
        subprocess.run(cmd, check=True, capture_output=True, timeout=300)
    except Exception:
        return None
    subs = sorted(out.glob("*.srt"))
    return subs[0] if subs else None


def ingest(source: str, store: SemanticStore, workdir: str = ".atlas_ingest") -> Doc:
    """Ingest a file path or video URL into semantic memory.

    Tutorial-style content additionally yields a draft skill playbook via
    the procedure extractor (returned in doc.meta["procedure_draft"]).
    """
    if source.startswith(("http://", "https://")):
        sub_path = fetch_video_transcript(source, workdir)
        text = (
            _text(sub_path)
            if sub_path
            else f"[no subtitles found for {source}; queue Whisper ASR]"
        )
    else:
        path = Path(source)
        handler = _HANDLERS.get(path.suffix.lower())
        text = handler(path) if handler else f"[unsupported format {path.suffix}]"

    doc = Doc(id=uuid.uuid4().hex[:12], text=text, source=source)
    procedure = extract_procedure(text)
    if procedure:
        doc.meta["procedure_draft"] = procedure
    store.add(doc)
    return doc
